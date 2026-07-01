#!/usr/bin/env python3
"""Turn a document into a presenter deck — cross-platform, permissively licensed.

Two modes:
  SHOW      render the document's real pages/slides to images (keep the original design).
  GENERATE  parse the document into sections + pull out its images, so an agent can author a
            meaningful deck (outline -> bullets + spoken narration, with the source's figures).

Default mode by input:
  .pptx/.ppt  -> show     (it's already a designed deck)
  .docx/.doc  -> generate (it's prose, not slides)
  .pdf        -> auto: a slide-shaped PDF (landscape, sparse) shows; a document (portrait,
                 dense) generates. Override with mode="show"/"generate".

Rendering engine: PDF -> images via **pypdfium2** (Chrome's PDFium; Apache/BSD; the wheel bundles
the native binary, so it installs by pip on Windows/macOS/Linux with NO system dependency).
PowerPoint/Word are rendered by first converting to PDF — via Microsoft Office (Windows COM, the
crispest path) or LibreOffice (cross-platform) — then pypdfium2. If neither is available, SHOW
degrades to GENERATE (text + extracted images) so you still get a usable deck, never a hard stop.

Image extraction uses **pikepdf** (PDF, lossless) and the Office file's own media (zip).
Permissive deps only: pypdfium2, pikepdf, pillow, python-pptx, python-docx (+ optional pywin32).
MIT. https://agentcall.dev
"""
import argparse
import glob
import json
import os
import shutil
import subprocess
import sys
import tempfile
import warnings
import zipfile

# UTF-8 out on every platform so status glyphs (→ etc.) never crash a Windows cp1252 console.
for _stream in ("stdout", "stderr"):
    try:
        getattr(sys, _stream).reconfigure(encoding="utf-8", errors="replace")
    except Exception:
        pass

_HERE = os.path.dirname(os.path.abspath(__file__))   # scripts/
_ROOT = os.path.dirname(_HERE)                       # the skill folder (scripts/ -> presenter/)
SUPPORTED = (".pdf", ".pptx", ".ppt", ".docx", ".doc")
_TARGET_W = 1600           # rendered page width in px (sharp on a video tile)
_DECK_MAX_WORDS = 60       # a PDF page with <= this many words looks slide-like (vs a document)


# ============================================================ document -> PDF (for Office files)

def _soffice():
    """Path to LibreOffice's soffice, if installed (cross-platform Office->PDF)."""
    for c in ("soffice", "libreoffice",
              r"C:\Program Files\LibreOffice\program\soffice.exe",
              r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
              "/Applications/LibreOffice.app/Contents/MacOS/soffice"):
        if os.path.isabs(c):
            if os.path.isfile(c):
                return c
        elif shutil.which(c):
            return shutil.which(c)
    return None


def _libreoffice_to_pdf(src, out_dir):
    so = _soffice()
    if not so:
        return None
    # a private profile dir avoids the single-instance lock when runs overlap
    profile = os.path.join(out_dir, "_lo_profile")
    # file:///C:/... on Windows, file:///tmp/... on POSIX (lstrip stops a malformed 4-slash URI)
    profile_uri = "file:///" + os.path.abspath(profile).replace("\\", "/").lstrip("/")
    try:
        subprocess.run(
            [so, "-env:UserInstallation=" + profile_uri,
             "--headless", "--convert-to", "pdf", "--outdir", out_dir, os.path.abspath(src)],
            check=True, timeout=180, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    except (subprocess.CalledProcessError, subprocess.TimeoutExpired, OSError) as e:
        # LibreOffice being INSTALLED but failing (corrupt file, flake, huge doc) must not hard-stop
        # the conversion — say so and let the caller's fallback chain take over.
        print(f"  LibreOffice couldn't convert {os.path.basename(src)} ({e}) — trying the fallback path.",
              file=sys.stderr)
        return None
    out = os.path.join(out_dir, os.path.splitext(os.path.basename(src))[0] + ".pdf")
    return out if os.path.isfile(out) else None


def _office_to_pdf(src, ext, out_dir):
    """Use installed Microsoft Office (Windows COM) to save the doc as PDF. Returns path or None."""
    if sys.platform != "win32":
        return None
    try:
        import win32com.client as win32
    except Exception:
        return None
    src = os.path.abspath(src)
    pdf = os.path.join(out_dir, os.path.splitext(os.path.basename(src))[0] + ".pdf")
    try:
        if ext in (".docx", ".doc"):
            # DispatchEx = a PRIVATE Word instance, so we never Quit() the user's own open Word.
            app = win32.DispatchEx("Word.Application")
            doc = None
            try:
                doc = app.Documents.Open(src, ReadOnly=1)
                doc.SaveAs(pdf, FileFormat=17)          # wdFormatPDF
            finally:
                try:
                    if doc is not None:
                        doc.Close(False)
                except Exception:
                    pass
                try:
                    app.Quit()
                except Exception:
                    pass
        elif ext in (".pptx", ".ppt"):
            # PowerPoint is single-instance: Dispatch attaches to the USER'S PowerPoint if it's open.
            # Only Quit() when nothing else is open — never tear down the user's own session (they may
            # literally have this deck open while asking the bot to present it).
            app = win32.Dispatch("PowerPoint.Application")
            pres = None
            try:
                pres = app.Presentations.Open(src, ReadOnly=1, WithWindow=False)
                pres.SaveAs(pdf, 32)                    # ppSaveAsPDF
            finally:
                try:
                    if pres is not None:
                        pres.Close()
                except Exception:
                    pass
                try:
                    if app.Presentations.Count == 0:
                        app.Quit()
                except Exception:
                    pass
        else:
            return None
    except Exception:
        return None
    return pdf if os.path.isfile(pdf) else None


def _office_to_pdf_any(src, ext, work):
    """Office -> PDF via MS Office (fast, Windows) then LibreOffice (cross-platform)."""
    return _office_to_pdf(src, ext, work) or _libreoffice_to_pdf(src, work)


# ============================================================ PDF -> page images (pypdfium2)

def _render_pdf_pages(pdf_path, img_dir):
    """Render each PDF page to slideN.png and grab its text. Returns slide dicts (show mode)."""
    import pypdfium2 as pdfium
    os.makedirs(img_dir, exist_ok=True)
    for f in glob.glob(os.path.join(img_dir, "slide*.png")):
        os.remove(f)
    pdf = pdfium.PdfDocument(pdf_path)
    slides = []
    try:
        for i in range(len(pdf)):
            page = pdf[i]
            w_pt = page.get_size()[0]
            scale = _TARGET_W / max(1.0, w_pt)          # default scale=1 is 72dpi (blurry) — scale up
            name = f"slide{i + 1}.png"
            page.render(scale=scale, draw_annots=True).to_pil().convert("RGB").save(os.path.join(img_dir, name))
            try:
                text = page.get_textpage().get_text_bounded() or ""
            except Exception:
                text = ""
            slides.append({"image": name, "text": text.strip(), "notes": ""})
            page.close()
    finally:
        pdf.close()
    return slides


def _pdf_page_geometry(pdf_path):
    """(avg_aspect_w_over_h, avg_words_per_page) — to tell a slide deck from a document."""
    import pypdfium2 as pdfium
    pdf = pdfium.PdfDocument(pdf_path)
    try:
        n = len(pdf)
        if not n:
            return 1.0, 0
        ar, words = 0.0, 0
        sample = min(n, 5)
        for i in range(sample):
            page = pdf[i]
            w, h = page.get_size()
            ar += (w / h) if h else 1.0
            try:
                words += len((page.get_textpage().get_text_bounded() or "").split())
            except Exception:
                pass
            page.close()
        return ar / sample, words / sample
    finally:
        pdf.close()


# ============================================================ image extraction (generate mode)

def _save_pil_rgb(pil, path):
    try:
        if pil.mode not in ("RGB", "L"):
            pil = pil.convert("RGB")
        pil.save(path)
        return True
    except Exception:
        return False


def _extract_pdf_images(pdf_path, img_dir, min_px=120):
    """Pull embedded raster images out of a PDF (pikepdf). Returns [(filename, page_index)]."""
    import pikepdf
    os.makedirs(img_dir, exist_ok=True)
    found = []
    try:
        pdf = pikepdf.open(pdf_path)
    except Exception:
        return found
    try:
        idx = 0
        for pageno, page in enumerate(pdf.pages):
            # pikepdf 9 deprecated Page.images (its replacement has a different return shape across
            # versions); we keep the stable .images access and just silence the cosmetic warning so
            # conversion output stays clean on every pikepdf version.
            with warnings.catch_warnings():
                warnings.simplefilter("ignore", DeprecationWarning)
                images = getattr(page, "images", {}) or {}
            for _name, raw in images.items():
                try:
                    pim = pikepdf.PdfImage(raw)
                    pil = pim.as_pil_image()
                except Exception:
                    continue
                if min(pil.size) < min_px:        # skip icons/bullets/rules
                    continue
                idx += 1
                fn = f"fig{idx}.png"
                if _save_pil_rgb(pil, os.path.join(img_dir, fn)):
                    found.append((fn, pageno))
    finally:
        pdf.close()
    return found


def _zip_media(path, img_dir, prefix="media"):
    """Extract embedded images from an Office (zip) file's media folder. Returns [filename]."""
    os.makedirs(img_dir, exist_ok=True)
    out = []
    try:
        z = zipfile.ZipFile(path)
    except Exception:
        return out
    with z:
        i = 0
        for n in z.namelist():
            low = n.lower()
            if ("/media/" in low) and low.rsplit(".", 1)[-1] in ("png", "jpg", "jpeg", "gif", "bmp", "tiff", "emf", "wmf"):
                ext = low.rsplit(".", 1)[-1]
                if ext in ("emf", "wmf"):           # vector metafiles a browser can't show — skip
                    continue
                i += 1
                fn = f"{prefix}{i}.{ 'jpg' if ext=='jpeg' else ext }"
                try:
                    with z.open(n) as src, open(os.path.join(img_dir, fn), "wb") as dst:
                        shutil.copyfileobj(src, dst)
                    out.append(fn)
                except Exception:
                    pass
    return out


# ============================================================ parse source (generate intermediate)

def _clean_lines(text):
    return [ln.strip() for ln in (text or "").splitlines() if ln.strip()]


def _parse_docx(path, img_dir):
    """Word -> sections [{title, text, images[]}] split by heading, with each embedded
    image attached to the section it appears in (so a figure lands on the right slide)."""
    from docx import Document
    from docx.oxml.ns import qn
    os.makedirs(img_dir, exist_ok=True)
    doc = Document(path)
    sections = []
    cur = {"title": "", "text": [], "images": []}
    n = [0]

    def save_blip(rid):
        try:
            part = doc.part.related_parts[rid]
            ext = os.path.splitext(part.partname)[1].lower().lstrip(".") or "png"
            if ext in ("emf", "wmf"):           # vector metafiles a browser can't show
                return None
            n[0] += 1
            fn = f"fig{n[0]}.{'jpg' if ext == 'jpeg' else ext}"
            with open(os.path.join(img_dir, fn), "wb") as f:
                f.write(part.blob)
            return fn
        except Exception:
            return None

    for p in doc.paragraphs:
        for blip in p._p.findall(".//" + qn("a:blip")):       # images in this paragraph, in order
            rid = blip.get(qn("r:embed"))
            if rid:
                fn = save_blip(rid)
                if fn:
                    cur["images"].append(fn)
        t = (p.text or "").strip()
        style = (p.style.name or "").lower() if p.style else ""
        if style.startswith("heading") or style == "title":
            if cur["title"] or cur["text"] or cur["images"]:
                sections.append(cur)
            cur = {"title": t, "text": [], "images": []}
        elif t:
            cur["text"].append(t)
    if cur["title"] or cur["text"] or cur["images"]:
        sections.append(cur)
    for s in sections:
        s["text"] = "\n".join(s["text"])
    all_images = [im for s in sections for im in s["images"]]
    return {"sections": sections, "all_images": all_images}


def _parse_pptx(path, img_dir):
    """PowerPoint -> per-slide {title, text, notes, images} (used only when show isn't possible)."""
    from pptx import Presentation
    prs = Presentation(path)
    media = _zip_media(path, img_dir, "media")
    sections = []
    for s in prs.slides:
        title = ""
        try:
            if s.shapes.title is not None:
                title = (s.shapes.title.text or "").strip()
        except Exception:
            pass
        lines = []
        for sh in s.shapes:
            if sh.has_text_frame:
                tx = (sh.text_frame.text or "").strip()
                if tx and tx != title:
                    lines.append(tx)
        notes = ""
        if s.has_notes_slide and s.notes_slide.notes_text_frame:
            notes = (s.notes_slide.notes_text_frame.text or "").strip()
        sections.append({"title": title, "text": "\n".join(lines), "notes": notes, "images": []})
    return {"sections": sections, "all_images": media}


def _pptx_notes(path):
    """Per-slide speaker notes from a .pptx, in slide order. Lets SHOW mode narrate a PowerPoint
    with NO agent/LLM — the bot reads the notes you already wrote. Returns [] on any failure."""
    try:
        from pptx import Presentation
        prs = Presentation(path)
    except Exception:
        return []
    out = []
    for s in prs.slides:
        nt = ""
        try:
            if s.has_notes_slide and s.notes_slide.notes_text_frame:
                nt = (s.notes_slide.notes_text_frame.text or "").strip()
        except Exception:
            nt = ""
        out.append(nt)
    return out


def _parse_pdf(path, img_dir):
    """PDF -> per-page {text} sections + extracted figures (attached to their page)."""
    import pypdfium2 as pdfium
    figs = _extract_pdf_images(path, img_dir)
    by_page = {}
    for fn, pageno in figs:
        by_page.setdefault(pageno, []).append(fn)
    pdf = pdfium.PdfDocument(path)
    sections = []
    try:
        for i in range(len(pdf)):
            page = pdf[i]
            try:
                text = (page.get_textpage().get_text_bounded() or "").strip()
            except Exception:
                text = ""
            page.close()
            sections.append({"title": "", "text": text, "images": by_page.get(i, [])})
    finally:
        pdf.close()
    return {"sections": sections, "all_images": [f for f, _ in figs]}


def parse_source(path, ext, img_dir):
    if ext in (".doc", ".ppt"):
        # Legacy OLE2 binaries — python-docx/python-pptx are OOXML-only and CRASH on these. Route
        # through Office/LibreOffice -> PDF -> parse; with no converter, fail with a clear message
        # instead of a cryptic "Package not found".
        with tempfile.TemporaryDirectory() as work:
            pdf = _office_to_pdf_any(path, ext, work)
            if not pdf:
                raise RuntimeError(
                    f"'{os.path.basename(path)}' is a legacy binary Office format — install "
                    "Word/PowerPoint or LibreOffice to convert it, or save it as "
                    f".{ 'docx' if ext == '.doc' else 'pptx' } first.")
            return _parse_pdf(pdf, img_dir)
    if ext == ".docx":
        return _parse_docx(path, img_dir)
    if ext == ".pptx":
        return _parse_pptx(path, img_dir)
    return _parse_pdf(path, img_dir)


# ============================================================ deck building

def _mechanical_deck(parsed):
    """A draft deck (titles + bullets) straight from parsed sections — NO narration invented. `notes`
    carries the section's REAL speaker notes only (usually none for a document); it is NEVER the raw
    body text or the heading, because reading those aloud is word-reading, not presenting. The agent
    authors real `notes` from source.json; until then the deck is flagged `needs_narration` (by
    convert) and present.py refuses to present it. Contentless sections (e.g. scanned/image-only PDF
    pages that yielded nothing) are SKIPPED — an all-empty parse returns [] so convert raises clearly."""
    slides = []
    for sec in parsed["sections"]:
        body = _clean_lines(sec.get("text", ""))
        head = sec.get("title", "") or (body[0] if body else "")
        bullets = [b for b in body if b != head][:6]
        notes = (sec.get("notes", "") or "").strip()   # real speaker notes only — never scraped text
        img = (sec.get("images") or [None])[0]
        if not (head or bullets or img):
            continue                                   # nothing to show OR say — not a slide
        slide = {"title": head, "bullets": bullets, "notes": notes}
        if img:
            slide["image"] = img
        slides.append(slide)
    return slides


def _resolve_mode(ext, mode, path):
    if mode in ("show", "generate"):
        return mode
    # auto
    if ext in (".pptx", ".ppt"):
        return "show"
    if ext in (".docx", ".doc"):
        return "generate"
    # PDF: slide-shaped (landscape + sparse) -> show; otherwise a document -> generate
    try:
        ar, words = _pdf_page_geometry(path)
        if ar >= 1.2 and words <= _DECK_MAX_WORDS:
            return "show"
    except Exception:
        pass
    return "generate"


def convert(path, mode="auto", out_dir=None):
    """Convert a document into a deck dir. Returns the deck.json path.

    mode: "show" (render real pages/slides), "generate" (parse for an agent to author),
          or "auto" (decide by file type / PDF shape).
    In generate mode a `source.json` intermediate is also written for the agent to author from.
    """
    path = os.path.abspath(path)
    if not os.path.isfile(path):
        raise FileNotFoundError(path)
    ext = os.path.splitext(path)[1].lower()
    if ext not in SUPPORTED:
        raise ValueError(f"Unsupported file type '{ext}'. Supported: {', '.join(SUPPORTED)}")

    base = os.path.splitext(os.path.basename(path))[0]
    if out_dir is None:
        out_dir = os.path.join(_ROOT, "decks", base)
    img_dir = os.path.join(out_dir, "img")

    # NEVER clobber authored work. If this document was converted before AND every slide's narration
    # has been authored (agent or human) AND the source file hasn't changed since, reuse that deck:
    # repeat presents are instant (no re-render) and the authored `notes` survive. If the SOURCE is
    # newer, the content changed — re-convert (and the needs_narration guard forces re-authoring).
    existing_path = os.path.join(out_dir, "deck.json")
    if os.path.isfile(existing_path) and os.path.getmtime(existing_path) >= os.path.getmtime(path):
        try:
            existing = json.loads(open(existing_path, encoding="utf-8-sig").read())
            sl = existing.get("slides") if isinstance(existing, dict) else None
            if sl and all((s.get("notes") or s.get("say") or "").strip() for s in sl):
                print(f"  Reusing the authored deck for {os.path.basename(path)} (source unchanged): {existing_path}")
                print("  (edit the document, or delete that deck folder, to force a fresh conversion)")
                return existing_path
        except (OSError, ValueError):
            pass                                   # unreadable existing deck → just re-convert

    os.makedirs(img_dir, exist_ok=True)
    deck_title = base.replace("_", " ").replace("-", " ").title()
    chosen = _resolve_mode(ext, mode, path)

    slides = None
    source = None
    with tempfile.TemporaryDirectory() as work:
        if chosen == "show":
            if ext == ".pdf":
                slides = _render_pdf_pages(path, img_dir)
            else:
                pdf = _office_to_pdf_any(path, ext, work)
                if pdf:
                    slides = _render_pdf_pages(pdf, img_dir)
                    # Narrate a PowerPoint with NO agent: attach each slide's speaker notes (if the
                    # rendered page count matches, so hidden slides can't misalign the narration).
                    if ext == ".pptx":
                        notes = _pptx_notes(path)
                        if notes and len(notes) == len(slides):
                            for sl, nt in zip(slides, notes):
                                sl["notes"] = nt
                        elif any(n.strip() for n in notes):
                            print(f"  Found speaker notes for {len(notes)} slides but rendered "
                                  f"{len(slides)} pages (hidden slides?) — notes NOT attached to avoid "
                                  "misalignment; author them via source.json.", file=sys.stderr)
                else:
                    print("  No PowerPoint/LibreOffice found to render the slides exactly — "
                          "falling back to a text+image deck. Install LibreOffice for exact-design "
                          "rendering.", file=sys.stderr)
                    chosen = "generate"
        if chosen == "generate":
            source = parse_source(path, ext, img_dir)
            slides = _mechanical_deck(source)

    if not slides:
        raise RuntimeError(f"No slides/text could be extracted from {os.path.basename(path)} — "
                           "is it a scanned or image-only file? Convert it to a real PDF/PPTX, "
                           "or author a deck.json by hand.")

    # Source unchanged but only PARTIALLY authored (the fully-authored case returned early above):
    # carry authored `notes` over by slide index, so a re-convert can never destroy narration work
    # in progress. Only safe when the source hasn't changed (same slides) and counts match.
    if os.path.isfile(existing_path) and os.path.getmtime(existing_path) >= os.path.getmtime(path):
        try:
            old = json.loads(open(existing_path, encoding="utf-8-sig").read())
            old_sl = old.get("slides") if isinstance(old, dict) else None
            if old_sl and len(old_sl) == len(slides):
                for new, prev in zip(slides, old_sl):
                    if not (new.get("notes") or "").strip() and (prev.get("notes") or "").strip():
                        new["notes"] = prev["notes"].strip()
        except (OSError, ValueError):
            pass

    # A CONVERTED deck where ANY slide lacks human speaker notes has slides with nothing meaningful
    # to SAY: on-screen text/bullets are a layout, not narration. We do NOT invent narration and we do
    # NOT read the layout aloud — either makes the bot a word-reader (a partially-noted deck would
    # otherwise silently flip through its un-noted slides). When any slide is unnarrated we flag
    # `needs_narration` and write source.json; the agent MUST author the missing `notes` (present.py
    # refuses until every slide has them). Covers show AND generate paths.
    unnarrated = not all((sl.get("notes") or "").strip() for sl in slides)

    # source.json — the reference an agent authors narration from (extracted text/sections + images).
    if source is not None or unnarrated:
        src_obj = {"title": deck_title, "source": path}
        if source is not None:
            src_obj.update(source)                         # generate: parsed sections + figures
        else:                                              # show-without-notes: rendered slides + text
            src_obj["slides"] = [{"n": i + 1, "image": sl.get("image", ""),
                                  "text": (sl.get("text") or "").strip()}
                                 for i, sl in enumerate(slides)]
        with open(os.path.join(out_dir, "source.json"), "w", encoding="utf-8") as f:
            json.dump(src_obj, f, indent=2, ensure_ascii=False)

    for sl in slides:
        sl.setdefault("title", "")
        sl["notes"] = (sl.get("notes") or "").strip()      # human speaker notes only — never scraped text
        sl.pop("text", None)                               # `text` is never a run-time narration source

    deck = {"title": deck_title, "mode": chosen, "slides": slides}
    if unnarrated:
        deck["needs_narration"] = True

    deck_path = os.path.join(out_dir, "deck.json")
    with open(deck_path, "w", encoding="utf-8") as f:
        json.dump(deck, f, indent=2, ensure_ascii=False)
    return deck_path


def main():
    ap = argparse.ArgumentParser(description="Convert a PDF / PowerPoint / Word file into a presenter deck.")
    ap.add_argument("doc", help="path to a .pdf, .pptx/.ppt, or .docx/.doc")
    ap.add_argument("--mode", choices=["auto", "show", "generate"], default="auto",
                    help="show real pages, generate a deck to author, or auto (default)")
    ap.add_argument("--out", default=None, help="output deck dir (default: decks/<name>/)")
    args = ap.parse_args()
    try:
        dp = convert(args.doc, mode=args.mode, out_dir=args.out)
    except Exception as e:
        print(f"Conversion failed: {e}")
        sys.exit(1)
    deck = json.load(open(dp, encoding="utf-8"))
    print(f"Deck written: {dp}  ({len(deck['slides'])} slides, mode={deck.get('mode')})")
    if deck.get("mode") == "generate":
        print("  Also wrote source.json — an agent should read it and author meaningful")
        print("  narration/bullets (and attach the extracted figures) into deck.json.")
    elif deck.get("needs_narration"):
        print("  NOTE: these real slides have NO speaker notes, so `notes` is empty. Reading the")
        print("  on-screen text aloud would make the bot a word-reader, not a presenter. An agent")
        print("  should read source.json (and the slide images) and AUTHOR a spoken `notes` line")
        print("  per slide in deck.json before presenting. See SKILL.md → 'Presenting a document'.")


if __name__ == "__main__":
    main()
